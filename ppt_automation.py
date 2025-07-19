#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime


class PowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title, subtitle=None, author=None):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            
            if author:
                subtitle_shape.text += f"\n\nBy: {author}"
        
        return slide
    
    def add_content_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        content_shape = slide.placeholders[1]
        
        if isinstance(content, list):
            text_frame = content_shape.text_frame
            for i, item in enumerate(content):
                if i == 0:
                    text_frame.text = item
                else:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0
        else:
            content_shape.text = content
        
        return slide
    
    def add_image_slide(self, title, image_path, caption=None):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if os.path.exists(image_path):
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            slide.shapes.add_picture(image_path, left, top, width=width)
            
            if caption:
                left = Inches(1)
                top = Inches(6)
                width = Inches(8)
                height = Inches(1)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                p = text_frame.paragraphs[0]
                p.text = caption
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(12)
                p.font.italic = True
        
        return slide
    
    def add_table_slide(self, title, headers, data):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        rows = len(data) + 1
        cols = len(headers)
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
            paragraph = cell.text_frame.paragraphs[0]
            font = paragraph.font
            font.bold = True
            font.color.rgb = RGBColor(255, 255, 255)
        
        for row_idx, row_data in enumerate(data):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
        
        return slide
    
    def add_chart_placeholder_slide(self, title, description):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = f"Chart Placeholder: {description}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(100, 100, 100)
        
        return slide
    
    def add_section_divider(self, section_title):
        slide_layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(slide_layout)
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = section_title
                break
        
        return slide
    
    def save(self, filename=None):
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
        
        self.prs.save(filename)
        return filename


def create_sample_presentation():
    ppt = PowerPointGenerator()
    
    ppt.add_title_slide(
        "PowerPoint Automation with Python",
        "Creating presentations programmatically",
        "Python Developer"
    )
    
    ppt.add_content_slide(
        "Features",
        [
            "Create presentations programmatically",
            "Add various slide types",
            "Insert text, images, and tables",
            "Customize formatting",
            "Save with automatic naming"
        ]
    )
    
    ppt.add_section_divider("Content Examples")
    
    ppt.add_table_slide(
        "Sample Data Table",
        ["Product", "Q1 Sales", "Q2 Sales", "Q3 Sales"],
        [
            ["Product A", "$10,000", "$12,000", "$15,000"],
            ["Product B", "$8,000", "$9,500", "$11,000"],
            ["Product C", "$15,000", "$14,000", "$16,500"]
        ]
    )
    
    ppt.add_chart_placeholder_slide(
        "Sales Trends",
        "Quarterly sales data visualization"
    )
    
    ppt.add_content_slide(
        "Next Steps",
        [
            "Install python-pptx library",
            "Customize templates",
            "Add more slide types",
            "Integrate with data sources",
            "Automate report generation"
        ]
    )
    
    filename = ppt.save("sample_presentation.pptx")
    print(f"Presentation saved as: {filename}")


if __name__ == "__main__":
    create_sample_presentation()